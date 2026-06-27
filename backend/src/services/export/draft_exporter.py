"""DraftExporter (T250) — async PPTX export with source attribution."""

from __future__ import annotations

import io
import uuid
from datetime import datetime

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from src.core.observability import get_logger
from src.db.models import Draft, DraftExportJob, DraftSlide
from src.scheduler.queue import publish_ws_event
from src.services.export.source_attribution import add_source_to_slide
from src.storage.minio import put_object

logger = get_logger("export.draft_exporter")


class DraftExporter:
    def __init__(self, session: AsyncSession) -> None:
        self.session = session

    async def create_job(self, draft_id: uuid.UUID, owner_id: uuid.UUID) -> DraftExportJob:
        # Verify ownership
        draft = (
            await self.session.execute(
                select(Draft).where(Draft.id == draft_id, Draft.owner_id == owner_id)
            )
        ).scalar_one_or_none()
        if not draft:
            from src.core.errors import NotFoundError

            raise NotFoundError("Draft", str(draft_id))
        if not draft.slides:
            from src.core.errors import PPTagentError

            raise PPTagentError(
                code="PPTAGENT.DRAFT_EMPTY",
                message="Cannot export an empty draft",
                status_code=422,
            )
        job = DraftExportJob(
            draft_id=draft_id,
            status="queued",
            progress=0,
        )
        self.session.add(job)
        await self.session.flush()
        return job

    @staticmethod
    async def run_export(job_id: str) -> None:
        from src.db.session import get_session_factory

        factory = get_session_factory()
        async with factory() as session:
            job = (
                await session.execute(
                    select(DraftExportJob).where(DraftExportJob.id == uuid.UUID(job_id))
                )
            ).scalar_one_or_none()
            if not job:
                logger.error("export_job_not_found", job_id=job_id)
                return
            job.status = "running"
            job.started_at = datetime.utcnow()
            job.progress = 10
            await session.commit()

            try:
                pptx_bytes = await DraftExporter._build_pptx(session, job)
                key = f"drafts/{job.draft_id}/{job.id}.pptx"
                put_object(
                    bucket="ppt-hot",
                    key=key,
                    data=pptx_bytes,
                    content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
                job.pptx_path = f"s3://ppt-hot/{key}"
                job.status = "success"
                job.progress = 100
                job.finished_at = datetime.utcnow()
                await session.commit()
                await publish_ws_event(
                    f"draft:{job.draft_id}",
                    {
                        "type": "draft.exported",
                        "draft_id": str(job.draft_id),
                        "job_id": str(job.id),
                        "pptx_path": job.pptx_path,
                    },
                )
            except Exception as e:
                job.status = "failed"
                job.error_message = str(e)[:2000]
                job.finished_at = datetime.utcnow()
                await session.commit()
                logger.exception("export_failed", job_id=job_id, error=str(e))

    @staticmethod
    async def _build_pptx(session: AsyncSession, job: DraftExportJob) -> bytes:
        """Build a PPTX with one slide per DraftSlide + source attribution metadata."""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        slides = (
            await session.execute(
                select(DraftSlide)
                .where(DraftSlide.draft_id == job.draft_id)
                .order_by(DraftSlide.slide_order.asc())
            )
        ).scalars()

        for ds in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            if ds.title:
                slide.shapes.title.text = ds.title
            if ds.body_text:
                tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(4.5))
                tf = tx.text_frame
                tf.text = ds.body_text
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(14)
            # Add source attribution (R11) — written as a custom property on the slide
            add_source_to_slide(
                slide,
                source_type=ds.source_type.value,
                material_id=str(ds.material_id) if ds.material_id else None,
                stage_id=str(ds.generated_stage_id) if ds.generated_stage_id else None,
            )

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
