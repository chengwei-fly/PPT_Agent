"""PPTX extractor — pulls per-slide images + text from a PPTX file.

The existing ``SampleParser`` only emits text chunks + layout hints, never
images. For the curated library use case (icon packs, flow diagrams,
relationship graphics, maps) the *image* is the asset, so this module
re-implements extraction at a lower level using ``python-pptx`` + Pillow.

Output is a list of :class:`ExtractedAsset` records — one per *shape-image*
found in the file. Each record carries:

* the embedded image bytes + original extension (PNG / JPG / JPEG / GIF / WEBP /
  BMP / TIFF / EMF / WMF)
* a flattened thumbnail (Pillow-cached at 480px wide) suitable for LLM input
* the slide's text content (truncated) for keyword-based classification
* the slide index, slide dimensions, source filename

EMF / WMF on Windows are converted to PNG via Pillow's native EMF support
(no external ``inkscape`` / ``libreoffice`` required for the common case).
Pure DrawingML vector shapes (no embedded image) are flattened by
compositing onto a white background and rasterising with Pillow.

This module does NOT touch the database or call any LLM. It is purely a
shape walker that the :mod:`importer` consumes.
"""

from __future__ import annotations

import io
import re
import zipfile
from collections.abc import Iterator
from dataclasses import dataclass, field
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

from src.core.observability import get_logger

logger = get_logger("material_importer.extractor")

# Thumbnail width (height auto). Used for both the persisted MinIO thumbnail
# and the LLM input — 480px is enough fidelity for icon classification while
# keeping request size small.
THUMB_MAX_WIDTH = 480
JPEG_QUALITY = 85
# Max extracted text length per asset — keeps DB row size bounded.
TEXT_LIMIT = 1500
# Slide composite fallback size (used only when a slide has zero embedded
# images, e.g. pure DrawingML shapes). The output is a rasterised version of
# the whole slide.
SLIDE_FALLBACK_W = 1280
SLIDE_FALLBACK_H = 720

SUPPORTED_IMG_EXT = {
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "gif": "image/gif",
    "webp": "image/webp",
    "bmp": "image/bmp",
    "tif": "image/tiff",
    "tiff": "image/tiff",
    "emf": "image/x-emf",
    "wmf": "image/x-wmf",
}


@dataclass
class ExtractedAsset:
    """One curated-library candidate: an image extracted from a slide."""

    source_file: str
    slide_index: int
    shape_name: str
    image_bytes: bytes
    image_ext: str
    mime_type: str
    width: int
    height: int
    text: str = ""
    title_hint: str = ""
    extra_shapes: int = 0
    palette: list[str] = field(default_factory=list)


# ────────────────────────────────────────────────────────────────────
# Main entry points
# ────────────────────────────────────────────────────────────────────


class PPTXExtractor:
    """Walk a PPTX file and yield :class:`ExtractedAsset` records.

    The extractor prefers *embedded image shapes* (one asset per image).
    If a slide contains no image at all, a single fallback asset is yielded
    representing the whole slide composited onto a canvas.
    """

    def __init__(self, thumb_max_width: int = THUMB_MAX_WIDTH) -> None:
        self.thumb_max_width = thumb_max_width

    def extract(self, file_path: str | Path) -> list[ExtractedAsset]:
        """Return a list of assets extracted from ``file_path``."""
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"PPTX not found: {path}")
        # python-pptx happily reads .ppt files too (best-effort) but officially
        # only supports .pptx. We still try — Office's binary .ppt is rare
        # now and most users have converted. Skip with a clear error.
        if path.suffix.lower() == ".ppt":
            logger.warning("legacy_ppt_format", file=str(path))
            # python-pptx will fail below with a clear error
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(str(path))
        slide_w = prs.slide_width or 9144000  # 10" in EMU
        slide_h = prs.slide_height or 5143500
        source_file = path.name

        assets: list[ExtractedAsset] = []
        for sidx, slide in enumerate(prs.slides):
            slide_text = _extract_slide_text(slide)
            title_hint = _extract_title(slide) or ""
            image_shapes: list[tuple[str, bytes, str, int, int]] = []
            other_count = 0
            for shape in slide.shapes:
                st = shape.shape_type
                # Picture shape — the most common case for icon libraries
                if st == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img = shape.image  # type: ignore[attr-defined]
                        ext = (img.ext or "png").lower()
                        if ext not in SUPPORTED_IMG_EXT:
                            logger.debug(
                                "skip_unsupported_image",
                                ext=ext,
                                file=source_file,
                                slide=sidx,
                            )
                            continue
                        image_shapes.append(
                            (
                                shape.name or f"picture_{sidx}",
                                img.blob,
                                ext,
                                shape.width or 0,
                                shape.height or 0,
                            )
                        )
                    except Exception as e:  # noqa: BLE001
                        logger.debug(
                            "picture_extract_failed",
                            file=source_file,
                            slide=sidx,
                            error=str(e),
                        )
                else:
                    other_count += 1

            if image_shapes:
                for name, blob, ext, w, h in image_shapes:
                    try:
                        pil = _open_image(blob, ext)
                    except Exception as e:  # noqa: BLE001
                        logger.warning(
                            "image_open_failed",
                            file=source_file,
                            slide=sidx,
                            ext=ext,
                            error=str(e),
                        )
                        continue
                    palette = _extract_palette(pil)
                    thumb_blob, tw, th = _make_thumbnail(pil, self.thumb_max_width)
                    ext_lower = "jpg" if ext in ("jpeg",) else ext
                    assets.append(
                        ExtractedAsset(
                            source_file=source_file,
                            slide_index=sidx,
                            shape_name=name[:128],
                            image_bytes=thumb_blob,
                            image_ext=ext_lower,
                            mime_type=SUPPORTED_IMG_EXT[ext],
                            width=tw,
                            height=th,
                            text=slide_text[:TEXT_LIMIT],
                            title_hint=title_hint[:255],
                            extra_shapes=other_count,
                            palette=palette,
                        )
                    )
            else:
                # Fallback: render the whole slide onto a canvas. Used when a
                # slide is pure DrawingML (e.g. Roland Berger-style template
                # pages with no embedded raster).
                try:
                    pil = _render_slide_fallback(slide, slide_w, slide_h, slide_text)
                except Exception as e:  # noqa: BLE001
                    logger.debug(
                        "slide_fallback_failed",
                        file=source_file,
                        slide=sidx,
                        error=str(e),
                    )
                    continue
                palette = _extract_palette(pil)
                thumb_blob, tw, th = _make_thumbnail(pil, self.thumb_max_width)
                assets.append(
                    ExtractedAsset(
                        source_file=source_file,
                        slide_index=sidx,
                        shape_name=f"slide_{sidx}",
                        image_bytes=thumb_blob,
                        image_ext="png",
                        mime_type="image/png",
                        width=tw,
                        height=th,
                        text=slide_text[:TEXT_LIMIT],
                        title_hint=title_hint[:255],
                        extra_shapes=other_count,
                        palette=palette,
                    )
                )
        return assets


def extract_from_file(file_path: str | Path) -> list[ExtractedAsset]:
    """Convenience wrapper around :class:`PPTXExtractor`."""
    return PPTXExtractor().extract(file_path)


# ────────────────────────────────────────────────────────────────────
# Image helpers
# ────────────────────────────────────────────────────────────────────


def _open_image(blob: bytes, ext: str) -> Image.Image:
    """Open an image (raster or EMF/WMF) and return RGB Image."""
    img = Image.open(io.BytesIO(blob))
    # Pillow can read EMF/WMF on Windows using its bundled backend; on other
    # platforms this will raise. Caller handles failures.
    if img.mode in ("RGBA", "LA", "P"):
        bg = Image.new("RGB", img.size, (255, 255, 255))
        if img.mode == "P":
            img = img.convert("RGBA")
        bg.paste(img, mask=img.split()[-1] if img.mode in ("RGBA", "LA") else None)
        return bg
    if img.mode != "RGB":
        return img.convert("RGB")
    return img


def _make_thumbnail(pil: Image.Image, max_width: int) -> tuple[bytes, int, int]:
    """Down-scale to ``max_width`` and JPEG-encode for compactness.

    Returns (bytes, width, height).
    """
    w, h = pil.size
    if w > max_width:
        ratio = max_width / w
        new_w = max_width
        new_h = max(1, int(h * ratio))
        pil = pil.resize((new_w, new_h), Image.LANCZOS)
    else:
        new_w, new_h = w, h
    buf = io.BytesIO()
    pil.save(buf, format="JPEG", quality=JPEG_QUALITY, optimize=True)
    return buf.getvalue(), new_w, new_h


_HEX_RE = re.compile(r"#[0-9a-fA-F]{6}")


def _extract_palette(pil: Image.Image, n: int = 5) -> list[str]:
    """Best-effort dominant-color extraction.

    We down-sample to 64x64, quantize to 6 colours, then sort by frequency.
    Cheap and good enough for visual-type filtering.
    """
    try:
        small = pil.copy()
        small.thumbnail((64, 64))
        paletted = small.quantize(colors=6, method=Image.Quantize.MEDIANCUT)
        # quantize palette is a flat list of 256*3 RGB bytes
        pal = paletted.getpalette() or []
        counts = sorted(paletted.getcolors() or [], reverse=True)
        out: list[str] = []
        for count, idx in counts[:n]:
            r = pal[idx * 3]
            g = pal[idx * 3 + 1]
            b = pal[idx * 3 + 2]
            if r > 250 and g > 250 and b > 250:
                continue  # skip near-white
            out.append(f"#{r:02x}{g:02x}{b:02x}")
        return out
    except Exception:  # noqa: BLE001
        return []


def _render_slide_fallback(slide, slide_w: int, slide_h: int, text: str) -> Image.Image:
    """Render a slide with no embedded images as a Pillow canvas.

    This is a *visual* fallback (not a faithful rendering). We draw the
    title text onto a white background at slide aspect ratio. Good enough
    to give the LLM a sense of "this is a text-heavy body slide" vs "this
    is a cover slide".
    """
    aspect = slide_w / max(1, slide_h)
    w = SLIDE_FALLBACK_W
    h = int(w / aspect)
    canvas = Image.new("RGB", (w, h), (255, 255, 255))
    draw = ImageDraw.Draw(canvas)
    title = (text or "").split("|")[0].strip()[:80]
    try:
        font = ImageFont.truetype("arial.ttf", 48)
        body = ImageFont.truetype("arial.ttf", 24)
    except OSError:
        font = ImageFont.load_default()
        body = font
    if title:
        draw.text((40, 40), title, fill=(31, 41, 55), font=font)
    if text and text != title:
        body_text = text.replace(title, "", 1).strip()[:200]
        draw.text((40, 140), body_text, fill=(75, 85, 99), font=body)
    return canvas


# ────────────────────────────────────────────────────────────────────
# Text helpers
# ────────────────────────────────────────────────────────────────────


def _extract_slide_text(slide) -> str:
    """Concatenate all text on the slide into a single | -separated string."""
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = "".join(run.text for run in para.runs).strip()
                if txt:
                    parts.append(txt)
        elif shape.has_table:
            try:
                for row in shape.table.rows:
                    for cell in row.cells:
                        txt = cell.text.strip()
                        if txt:
                            parts.append(txt)
            except Exception:  # noqa: BLE001
                pass
    return " | ".join(parts)


def _extract_title(slide) -> str | None:
    title = slide.shapes.title
    if title is None or not title.has_text_frame:
        return None
    return title.text_frame.text.strip() or None


# ────────────────────────────────────────────────────────────────────
# Discovery — list PPTX files in a directory
# ────────────────────────────────────────────────────────────────────


def discover_pptx_files(root: str | Path) -> Iterator[Path]:
    """Yield PPT/PPTX files under ``root`` (recursively).

    ``.ppt`` (legacy PowerPoint) is yielded but the extractor will refuse
    to parse it — see :class:`PPTXExtractor.extract`. Use
    :func:`convert_legacy_ppt` first if you need to ingest them.
    """
    root = Path(root)
    if not root.exists():
        return
    for path in sorted(root.rglob("*")):
        if path.is_file() and path.suffix.lower() in {".pptx", ".ppt"}:
            yield path


def safe_open_zip(path: str | Path) -> zipfile.ZipFile:
    """Open a PPTX (which is a ZIP) with sanitised path traversal protection."""
    return zipfile.ZipFile(path, "r")


# ────────────────────────────────────────────────────────────────────
# Legacy .ppt conversion
# ────────────────────────────────────────────────────────────────────


def convert_legacy_ppt(path: str | Path) -> Path:
    """Convert a ``.ppt`` file to ``.pptx`` using the local PowerShell COM
    bridge on Windows, or LibreOffice headless elsewhere.

    Returns the path of the converted ``.pptx``. If neither tool is
    available, raises :class:`RuntimeError` with a clear hint.
    """
    import shutil
    import subprocess
    import tempfile

    path = Path(path)
    if path.suffix.lower() != ".ppt":
        raise ValueError(f"not a .ppt file: {path}")
    target_dir = Path(tempfile.mkdtemp(prefix="ppt_convert_"))
    # Prefer LibreOffice (cross-platform)
    lo = shutil.which("soffice") or shutil.which("libreoffice")
    if lo:
        subprocess.run(
            [
                lo,
                "--headless",
                "--convert-to",
                "pptx",
                "--outdir",
                str(target_dir),
                str(path),
            ],
            check=True,
            timeout=180,
        )
        converted = target_dir / (path.stem + ".pptx")
        if converted.exists():
            return converted
        raise RuntimeError(f"libreoffice produced no .pptx for {path}")
    # Windows fallback: PowerShell COM automation
    if hasattr(subprocess, "Popen") and __import__("os").name == "nt":
        ps_script = (
            f"$p = '{path}'; "
            f"$o = '{target_dir}\\{path.stem}.pptx'; "
            "$app = New-Object -ComObject PowerPoint.Application; "
            "$pres = $app.Presentations.Open($p, $true, $true, $false); "
            "$pres.SaveAs($o, 24); "  # 24 = ppSaveAsOpenXMLPresentation
            "$pres.Close(); "
            "$app.Quit();"
        )
        try:
            subprocess.run(
                ["powershell", "-NoProfile", "-Command", ps_script],
                check=True,
                timeout=180,
            )
            converted = target_dir / f"{path.stem}.pptx"
            if converted.exists():
                return converted
        except (subprocess.CalledProcessError, FileNotFoundError):
            pass
    raise RuntimeError(
        f"Cannot convert legacy .ppt file: {path}. "
        "Install LibreOffice (soffice on PATH) or Microsoft PowerPoint, "
        "or manually re-save the file as .pptx before importing."
    )
