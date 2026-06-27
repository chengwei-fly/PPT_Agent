"""SVG2PPTX tool — converts SVG slides to native PPTX format.

Workflow:

  1. Receive `slides: [{order, svg, title?}]` from the harness stage.
  2. Materialize each SVG to a temp file.
  3. Call `create_pptx_with_native_svg(...)` via the bridge module.
  4. Upload the resulting `.pptx` to MinIO and return the path.
  5. Clean up the temp dir.

A `python-pptx` fallback is kept as a defensive measure for the
case where the generation engine is unavailable (e.g. in CI).
"""

from __future__ import annotations

import asyncio
import io
import tempfile
import time
from pathlib import Path
from typing import Any

from src.core.observability import get_logger
from src.integrations.svg_pptx_bridge import create_pptx_with_native_svg
from src.storage.minio import put_object, result_bucket

logger = get_logger("svg2pptx")


class SVG2PPTXTool:
    name = "svg2pptx"
    description = (
        "Convert a list of slide SVG payloads into a single PPTX file. "
        "Uses the local SVG-to-PPTX generation engine via "
        "`src.integrations.svg_pptx_bridge`."
    )
    parameters: dict[str, Any] = {
        "type": "object",
        "properties": {
            "task_id": {"type": "string", "description": "Owning task UUID"},
            "slides": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "order": {"type": "integer"},
                        "svg": {"type": "string"},
                        "title": {"type": "string"},
                    },
                    "required": ["order", "svg"],
                },
                "description": "Slides in order",
            },
            "theme": {
                "type": "object",
                "description": "Theme overrides (colors / fonts / layout)",
                "properties": {
                    "palette": {"type": "array", "items": {"type": "string"}},
                    "font_family": {"type": "string"},
                    "layout": {"type": "string"},
                },
            },
            "notes": {
                "type": "object",
                "description": "Optional per-slide speaker notes (order → text).",
                "additionalProperties": {"type": "string"},
            },
        },
        "required": ["task_id", "slides"],
    }

    async def func(
        self,
        task_id: str,
        slides: list[dict[str, Any]],
        theme: dict | None = None,
        notes: dict[str, str] | None = None,
    ) -> dict[str, Any]:
        start = time.perf_counter()
        if not slides:
            raise ValueError("svg2pptx: slides must be non-empty")

        # Sort by `order` defensively.
        slides = sorted(slides, key=lambda s: s.get("order", 0))

        # Materialize SVGs to a temp dir (the local API wants
        # list[Path]). Use a per-call tempdir so concurrent tasks
        # don't collide.
        with tempfile.TemporaryDirectory(prefix=f"pptagent_svg_{task_id}_") as tmp:
            svg_dir = Path(tmp) / "svg"
            svg_dir.mkdir(parents=True, exist_ok=True)
            out_path = Path(tmp) / f"{task_id}.pptx"

            svg_paths: list[Path] = []
            for i, s in enumerate(slides):
                p = svg_dir / f"{i + 1:02d}.svg"
                p.write_text(s["svg"], encoding="utf-8")
                svg_paths.append(p)

            # Build the optional canvas_format from the theme hint
            canvas_format = (theme or {}).get("layout") or "169"
            use_native_shapes = bool((theme or {}).get("native_shapes", False))

            try:
                ok = await asyncio.to_thread(
                    create_pptx_with_native_svg,
                    svg_files=svg_paths,
                    output_path=out_path,
                    canvas_format=canvas_format,
                    verbose=False,
                    notes=notes,
                    enable_notes=bool(notes),
                    use_native_shapes=use_native_shapes,
                    use_compat_mode=True,
                )
            except Exception as e:
                logger.exception(
                    "svg_to_pptx_engine_failed",
                    task_id=task_id,
                    error=str(e),
                )
                return await self._fallback_pptx(
                    task_id,
                    slides,
                    theme,
                    start,
                    reason=str(e),
                )

            if not ok or not out_path.is_file():
                logger.error(
                    "svg_to_pptx_engine_returned_false",
                    task_id=task_id,
                    output_exists=out_path.is_file(),
                )
                return await self._fallback_pptx(
                    task_id,
                    slides,
                    theme,
                    start,
                    reason="create_pptx_with_native_svg returned False",
                )

            data = out_path.read_bytes()

        key = f"results/{task_id}.pptx"
        put_object(
            bucket=result_bucket(),
            key=key,
            data=data,
            content_type=(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
        )
        duration_ms = int((time.perf_counter() - start) * 1000)
        logger.info(
            "svg2pptx_engine_ok",
            task_id=task_id,
            slide_count=len(slides),
            bytes=len(data),
            duration_ms=duration_ms,
        )
        return {
            "pptx_path": f"s3://{result_bucket()}/{key}",
            "slide_count": len(slides),
            "bytes": len(data),
            "duration_ms": duration_ms,
            "engine": "svg-to-pptx",
        }

    async def _fallback_pptx(
        self,
        task_id: str,
        slides: list[dict[str, Any]],
        theme: dict | None,
        start: float,
        reason: str,
    ) -> dict[str, Any]:
        """Defensive fallback — should never be hit in production.

        Runs when the SVG-to-PPTX generation engine is unavailable
        (e.g. a CI sandbox). The orchestrator surfaces a clear
        warning when this path is taken.
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        for s in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
            if s.get("title"):
                slide.shapes.title.text = s["title"]
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(4.5))
            tf = tx.text_frame
            tf.text = f"[SVG: {len(s.get('svg', ''))} chars — fallback]"
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(10)

        buf = io.BytesIO()
        prs.save(buf)
        data = buf.getvalue()
        key = f"results/{task_id}.pptx"
        put_object(
            bucket=result_bucket(),
            key=key,
            data=data,
            content_type=(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
        )
        duration_ms = int((time.perf_counter() - start) * 1000)
        logger.warning(
            "using_python_pptx_fallback",
            task_id=task_id,
            slide_count=len(slides),
            duration_ms=duration_ms,
            reason=reason,
        )
        return {
            "pptx_path": f"s3://{result_bucket()}/{key}",
            "slide_count": len(slides),
            "bytes": len(data),
            "duration_ms": duration_ms,
            "engine": "python-pptx-fallback",
            "fallback_reason": reason,
        }
