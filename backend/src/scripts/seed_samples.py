#!/usr/bin/env python
"""Seed 5 typical PPTX samples into MinIO (T009).

Run after `uv run alembic upgrade head` and after MinIO buckets are created.
Uses pre-built fixtures from `backend/tests/fixtures/samples/`.
"""

from __future__ import annotations

import asyncio
import hashlib
import os
from pathlib import Path

from minio import Minio
from minio.error import S3Error

SAMPLES = [
    ("汇报-template.pptx", "report", "通用工作汇报模板（季度/年终）"),
    ("培训-template.pptx", "training", "内部培训课件模板（含章节、要点）"),
    ("方案-template.pptx", "proposal", "项目方案/招标书模板（封面+目录+方案+附录）"),
    ("数据-template.pptx", "data", "数据分析报告模板（图表+洞察）"),
    ("营销-template.pptx", "marketing", "市场活动/营销提案模板（视觉冲击）"),
]


def make_placeholder_pptx(target: Path) -> int:
    """Generate a minimal but valid PPTX fixture (1 slide, sample theme).

    For real fixtures, replace with actual PPTX samples in
    backend/tests/fixtures/samples/.
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        # Fall back to writing a minimal file so the script can run in CI
        # even if python-pptx isn't installed at runtime.
        target.write_bytes(b"PK\x03\x04placeholder-pptx-fixture")
        return 0

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    title = slide.shapes.title
    title.text = target.stem
    for para in title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(40)
    prs.save(str(target))
    return 1


async def main() -> None:
    fixtures_dir = Path(__file__).resolve().parent.parent.parent / "tests" / "fixtures" / "samples"
    fixtures_dir.mkdir(parents=True, exist_ok=True)

    client = Minio(
        endpoint=os.environ.get("S3_ENDPOINT", "localhost:9000"),
        access_key=os.environ.get("S3_ACCESS_KEY", "minioadmin"),
        secret_key=os.environ.get("S3_SECRET_KEY", "minioadmin"),
        secure=os.environ.get("S3_SECURE", "false") == "true",
    )
    bucket = os.environ.get("S3_BUCKET_HOT", "ppt-hot")
    if not client.bucket_exists(bucket):
        client.make_bucket(bucket)

    for filename, tag, desc in SAMPLES:
        local = fixtures_dir / filename
        if not local.exists():
            print(f"[seed] generating placeholder fixture: {filename}")
            make_placeholder_pptx(local)
        data = local.read_bytes()
        sha = hashlib.sha256(data).hexdigest()
        key = f"seed/{filename}"
        try:
            client.put_object(
                bucket_name=bucket,
                object_name=key,
                data=__import__("io").BytesIO(data),
                length=len(data),
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                metadata={"seed-tag": tag, "seed-description": desc, "seed-sha256": sha},
            )
            print(
                f"[seed] uploaded {filename} -> s3://{bucket}/{key}  ({len(data)} bytes, sha={sha[:12]})"
            )
        except S3Error as e:
            print(f"[seed] FAILED to upload {filename}: {e}")
            raise

    print(f"[seed] done. {len(SAMPLES)} samples seeded into '{bucket}'.")


if __name__ == "__main__":
    asyncio.run(main())
